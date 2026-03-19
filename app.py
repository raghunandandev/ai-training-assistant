# import streamlit as st
# import io
# from pptx import Presentation
# from processors import DocumentProcessor
# from ai_engine import AIEngine
# from tracker import ComplianceTracker

# # --- HELPER FUNCTION FOR PPTX GENERATION (Bonus 2) ---
# def create_ppt(training_data):
#     prs = Presentation()
    
#     # Title Slide
#     title_slide_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(title_slide_layout)
#     title = slide.shapes.title
#     subtitle = slide.placeholders[1]
#     title.text = training_data.get("title", "Training Module")
#     subtitle.text = "Generated AI Training Presentation"

#     # Summary Slide
#     bullet_slide_layout = prs.slide_layouts[1]
#     slide = prs.slides.add_slide(bullet_slide_layout)
#     shapes = slide.shapes
#     title_shape = shapes.title
#     body_shape = shapes.placeholders[1]
#     title_shape.text = "Executive Summary"
#     tf = body_shape.text_frame
#     tf.text = training_data.get("summary", "")

#     # Steps Slides
#     for i, step in enumerate(training_data.get("steps", [])):
#         slide = prs.slides.add_slide(bullet_slide_layout)
#         shapes = slide.shapes
#         title_shape = shapes.title
#         body_shape = shapes.placeholders[1]
#         title_shape.text = f"Step {i+1}"
#         tf = body_shape.text_frame
#         tf.text = step

#     # Save to memory stream
#     ppt_stream = io.BytesIO()
#     prs.save(ppt_stream)
#     ppt_stream.seek(0)
#     return ppt_stream

# # --- MAIN APP LOGIC ---
# def main():
#     st.set_page_config(page_title="AI Training System", layout="wide")
#     st.title("📄 SOP to AI Training System")
#     st.write("Upload a standard operating procedure to instantly generate a training module, slide deck, and compliance quiz.")

#     # Initialize objects
#     try:
#         ai = AIEngine()
#     except Exception as e:
#         st.error(f"Setup Error: {str(e)}")
#         return

#     tracker = ComplianceTracker()

#     # Session state to hold data so it doesn't disappear on button clicks
#     if "training_data" not in st.session_state:
#         st.session_state.training_data = None
#     if "quiz_submitted" not in st.session_state:
#         st.session_state.quiz_submitted = False

#     # 1. File Upload
#     uploaded_file = st.file_uploader("Upload SOP (PDF only)", type=["pdf"])

#     if uploaded_file is not None:
#         if st.button("Generate Training Materials", type="primary"):
#             with st.spinner("Extracting text and generating AI curriculum..."):
#                 doc = DocumentProcessor(uploaded_file)
#                 extracted_text = doc.extract_text()
                
#                 if extracted_text.startswith("Error"):
#                     st.error(extracted_text)
#                 else:
#                     data = ai.generate_training_module(extracted_text)
#                     if "error" in data:
#                         st.error(data["error"])
#                     else:
#                         st.session_state.training_data = data
#                         st.session_state.quiz_submitted = False
#                         st.rerun()

#     # 2. Display the UI if data exists
#     if st.session_state.training_data:
#         data = st.session_state.training_data
#         st.divider()
#         st.header(data.get("title", "Training Module"))
        
#         # Download PPTX Button
#         ppt_file = create_ppt(data)
#         st.download_button(
#             label="📥 Download Presentation (.pptx)",
#             data=ppt_file,
#             file_name="training_presentation.pptx",
#             mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
#         )

#         # Tabs for Content
#         tab1, tab2, tab3 = st.tabs(["📑 Summary", "⚙️ Step-by-Step", "🎓 Evaluation Quiz"])

#         with tab1:
#             st.subheader("Module Overview")
#             st.info(data.get("summary", "No summary available."))

#         with tab2:
#             st.subheader("Process Steps")
#             for i, step in enumerate(data.get("steps", [])):
#                 st.markdown(f"**{i+1}.** {step}")

#         with tab3:
#             st.subheader("Knowledge Check")
#             # Form for the quiz
#             with st.form("quiz_form"):
#                 # Removed required=True
#                 user_name = st.text_input("Employee Name:") 
                
#                 user_answers = {}
#                 for i, q in enumerate(data.get("quiz", [])):
#                     st.markdown(f"**Q{i+1}: {q['question']}**")
#                     user_answers[i] = st.radio("Select answer:", q['options'], key=f"q_{i}")
#                     st.write("")
                
#                 submitted = st.form_submit_button("Submit Answers")
                
#                 if submitted:
#                     # Manual check to ensure they entered a name
#                     if not user_name.strip():
#                         st.warning("⚠️ Please enter your Employee Name before submitting.")
#                     else:
#                         correct_count = 0
#                         for i, q in enumerate(data.get("quiz", [])):
#                             if user_answers[i] == q['answer']:
#                                 correct_count += 1
                        
#                         total = len(data.get("quiz", []))
#                         score_str = f"{correct_count}/{total}"
                        
#                         # Use tracker object to save data (Automation Bonus)
#                         tracker.log_score(user_name, data.get("title", "Unknown SOP"), score_str)
                        
#                         st.success(f"Quiz submitted! You scored {score_str}.")
#                         st.info("Your score has been automatically logged to the compliance database.")
            
#             # # Form for the quiz
#             # with st.form("quiz_form"):
#             #     user_name = st.text_input("Employee Name:", required=True)
                
#             #     user_answers = {}
#             #     for i, q in enumerate(data.get("quiz", [])):
#             #         st.markdown(f"**Q{i+1}: {q['question']}**")
#             #         user_answers[i] = st.radio("Select answer:", q['options'], key=f"q_{i}")
#             #         st.write("")
                
#             #     submitted = st.form_submit_button("Submit Answers")
                
#             #     if submitted:
#             #         correct_count = 0
#             #         for i, q in enumerate(data.get("quiz", [])):
#             #             if user_answers[i] == q['answer']:
#             #                 correct_count += 1
                    
#             #         total = len(data.get("quiz", []))
#             #         score_str = f"{correct_count}/{total}"
                    
#             #         # Use tracker object to save data (Automation Bonus)
#             #         tracker.log_score(user_name, data.get("title", "Unknown SOP"), score_str)
                    
#             #         st.success(f"Quiz submitted! You scored {score_str}.")
#             #         st.info("Your score has been automatically logged to the compliance database.")

# if __name__ == "__main__":
#     main()


#------------------------------>

import streamlit as st
import io
from pptx import Presentation
from processors import DocumentProcessor
from ai_engine import AIEngine
from tracker import ComplianceTracker

# # --- HELPER FUNCTION FOR PPTX GENERATION ---
# def create_ppt(training_data):
#     """Generates a slide deck based on the new Phase-Based JSON structure."""
#     prs = Presentation()
#     metadata = training_data.get("metadata", {})
    
#     # Slide 1: Title Slide
#     title_slide_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(title_slide_layout)
#     title = slide.shapes.title
#     subtitle = slide.placeholders[1]
#     title.text = metadata.get("title", "Training Module")
#     subtitle.text = f"Level: {metadata.get('complexity', 'Standard')} | Est. Time: {metadata.get('estimated_time_minutes', 'N/A')} mins"

#     # Slide 2: Executive Summary
#     bullet_slide_layout = prs.slide_layouts[1]
#     slide = prs.slides.add_slide(bullet_slide_layout)
#     shapes = slide.shapes
#     title_shape = shapes.title
#     body_shape = shapes.placeholders[1]
#     title_shape.text = "Executive Summary"
#     tf = body_shape.text_frame
#     tf.text = training_data.get("executive_summary", "Overview of the process.")

#     # Slides 3+: Action Phases
#     for phase in training_data.get("phases", []):
#         slide = prs.slides.add_slide(bullet_slide_layout)
#         shapes = slide.shapes
#         title_shape = shapes.title
#         body_shape = shapes.placeholders[1]
        
#         title_shape.text = f"Phase: {phase.get('phase_name', 'Step')}"
#         tf = body_shape.text_frame
        
#         for step in phase.get("steps", []):
#             p = tf.add_paragraph()
#             p.text = step
#             p.level = 0
            
#         if phase.get("critical_warning"):
#             p = tf.add_paragraph()
#             p.text = f"WARNING: {phase.get('critical_warning')}"
#             p.level = 1

#     # Save to memory stream
#     ppt_stream = io.BytesIO()
#     prs.save(ppt_stream)
#     ppt_stream.seek(0)
#     return ppt_stream

# --- HELPER FUNCTION FOR PPTX GENERATION ---
def create_ppt(training_data):
    """Generates a comprehensive, enterprise-ready slide deck."""
    prs = Presentation()
    metadata = training_data.get("metadata", {})
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = metadata.get("title", "Corporate Training Module")
    subtitle.text = f"Level: {metadata.get('complexity', 'Standard')} | Duration: {metadata.get('estimated_time_minutes', 'N/A')} mins\nGenerated by AI Training System"

    # 2. Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Executive Summary"
    tf = body_shape.text_frame
    tf.text = training_data.get("executive_summary", "Overview of the process.")

    # 3. Preparation & Tools Slide
    tools = metadata.get("tools_required", [])
    if tools:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        title_shape.text = "Preparation & Required Tools"
        tf = body_shape.text_frame
        for tool in tools:
            p = tf.add_paragraph()
            p.text = f"🔧 {tool}"
            p.level = 0

    # 4. Action Phases Slides
    for phase in training_data.get("phases", []):
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = f"Phase: {phase.get('phase_name', 'Process Step')}"
        tf = body_shape.text_frame
        
        for step in phase.get("steps", []):
            p = tf.add_paragraph()
            p.text = step
            p.level = 0
            
        if phase.get("critical_warning"):
            p = tf.add_paragraph()
            p.text = f"⚠️ CRITICAL WARNING: {phase.get('critical_warning')}"
            p.level = 1  # Indents the warning for visual emphasis

    # 5. Live Scenario Quiz Slides
    quiz_data = training_data.get("scenario_quiz", [])
    if quiz_data:
        # Add a transition/section slide
        section_layout = prs.slide_layouts[2] 
        slide = prs.slides.add_slide(section_layout)
        slide.shapes.title.text = "Knowledge Check"
        slide.placeholders[1].text = "Group Scenario Evaluation"

        # Add one slide per scenario
        for i, q in enumerate(quiz_data):
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = f"Scenario {i+1}"
            tf = body_shape.text_frame
            
            # The Scenario
            p = tf.add_paragraph()
            p.text = q.get("scenario", "")
            p.level = 0
            
            # The Question
            p = tf.add_paragraph()
            p.text = f"Question: {q.get('question', '')}"
            p.level = 1
            
            # The Options
            for opt in q.get("options", []):
                p = tf.add_paragraph()
                p.text = f"• {opt}"
                p.level = 2

    # Save to memory stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# --- MAIN APP LOGIC ---
def main():
    st.set_page_config(page_title="AI Training System", layout="wide", page_icon="🏢")
    st.title("📄 Enterprise SOP Training Generator")
    st.write("Transform raw standard operating procedures into scenario-based training modules.")

    # Initialize objects
    try:
        ai = AIEngine()
    except Exception as e:
        st.error(f"Setup Error: {str(e)}")
        return

    tracker = ComplianceTracker()

    # Session state
    if "training_data" not in st.session_state:
        st.session_state.training_data = None
    if "quiz_submitted" not in st.session_state:
        st.session_state.quiz_submitted = False

    # 1. File Upload
    uploaded_file = st.file_uploader("Upload SOP (PDF only)", type=["pdf"])

    if uploaded_file is not None:
        if st.button("Generate Training Materials", type="primary"):
            with st.spinner("Analyzing document and engineering curriculum..."):
                doc = DocumentProcessor(uploaded_file)
                extracted_text = doc.extract_text()
                
                if extracted_text.startswith("Error"):
                    st.error(extracted_text)
                else:
                    data = ai.generate_training_module(extracted_text)
                    if "error" in data:
                        st.error(data["error"])
                    else:
                        st.session_state.training_data = data
                        st.session_state.quiz_submitted = False
                        st.rerun()

    # 2. Display the Dashboard UI
    if st.session_state.training_data:
        data = st.session_state.training_data
        metadata = data.get("metadata", {})
        
        st.divider()
        st.header(metadata.get("title", "Training Module"))
        
        # --- DASHBOARD METRICS ---
        col1, col2, col3, col4 = st.columns(4)
        col1.metric("⏱️ Est. Time", f"{metadata.get('estimated_time_minutes', '--')} mins")
        col2.metric("📊 Complexity", metadata.get("complexity", "--"))
        col3.metric("🛠️ Tools Required", len(metadata.get("tools_required", [])))
        
        with col4:
            st.write("") # Spacing alignment
            ppt_file = create_ppt(data)
            st.download_button(
                label="📥 Download Deck (.pptx)",
                data=ppt_file,
                file_name="training_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

        st.write("") # Spacing

        # --- TABS ---
        tab1, tab2, tab3 = st.tabs(["📑 Overview & Prep", "⚙️ Action Phases", "🎓 Scenario Eval"])

        # TAB 1: OVERVIEW
        with tab1:
            st.subheader("Executive Summary")
            st.info(data.get("executive_summary", "No summary available."))
            
            st.subheader("Required Tools & Access")
            tools = metadata.get("tools_required", [])
            if tools:
                for tool in tools:
                    st.markdown(f"- {tool}")
            else:
                st.write("No specific tools required for this process.")

        # TAB 2: ACTION PHASES
        with tab2:
            st.subheader("Process Execution")
            for phase in data.get("phases", []):
                with st.expander(f"Phase: {phase.get('phase_name', 'Unnamed Phase')}", expanded=True):
                    for i, step in enumerate(phase.get("steps", [])):
                        st.markdown(f"**{i+1}.** {step}")
                    
                    if phase.get("critical_warning"):
                        st.warning(f"⚠️ **CRITICAL RISK:** {phase.get('critical_warning')}")

        # TAB 3: SCENARIO QUIZ
        with tab3:
            st.subheader("Scenario-Based Knowledge Check")
            
            with st.form("quiz_form"):
                user_name = st.text_input("Employee Name:")
                
                user_answers = {}
                for i, q in enumerate(data.get("scenario_quiz", [])):
                    st.markdown(f"### Scenario {i+1}")
                    st.write(q.get("scenario", ""))
                    st.markdown(f"**Question:** {q.get('question', '')}")
                    
                    user_answers[i] = st.radio("Select best action:", q.get('options', []), key=f"q_{i}")
                    st.divider()
                
                submitted = st.form_submit_button("Submit Evaluation")
                
                if submitted:
                    if not user_name.strip():
                        st.warning("⚠️ Please enter your Employee Name before submitting.")
                    else:
                        correct_count = 0
                        for i, q in enumerate(data.get("scenario_quiz", [])):
                            if user_answers[i] == q.get('answer'):
                                correct_count += 1
                                st.success(f"**Scenario {i+1}: Correct!** {q.get('explanation', '')}")
                            else:
                                st.error(f"**Scenario {i+1}: Incorrect.** The right answer was '{q.get('answer')}'. {q.get('explanation', '')}")
                        
                        total = len(data.get("scenario_quiz", []))
                        score_str = f"{correct_count}/{total}"
                        
                        # Automation Database Log
                        tracker.log_score(user_name, metadata.get("title", "Unknown SOP"), score_str)
                        
                        st.info(f"Final Score: {score_str}. This evaluation has been permanently logged in the compliance database.")

if __name__ == "__main__":
    main()